import asyncio
import csv
import io
import logging
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


class ParsedDocument:
    def __init__(self, content: str, metadata: dict[str, Any]):
        self.content = content
        self.metadata = metadata


class BaseParser:
    def __init__(self, timeout: int = 30):
        self.timeout = timeout

    async def parse(self, content: bytes) -> ParsedDocument:
        raise NotImplementedError


class PDFParser(BaseParser):
    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            return await asyncio.get_event_loop().run_in_executor(None, self._parse_sync, content)
        except Exception as e:
            logger.error(f"PDF parsing failed: {e}")
            return ParsedDocument(f"PDF parsing error: {e}", {"format": "pdf", "error": str(e)})

    def _parse_sync(self, content: bytes) -> ParsedDocument:
        text_parts = []
        num_pages = 0

        try:
            import fitz
            doc = fitz.open(stream=content, filetype="pdf")
            num_pages = len(doc)
            for page_num in range(num_pages):
                page = doc[page_num]
                text_parts.append(page.get_text())
            doc.close()
            logger.info(f"Extracted {num_pages} pages via PyMuPDF")
        except ImportError:
            try:
                import pdfplumber
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    num_pages = len(pdf.pages)
                    for page in pdf.pages:
                        page_text = page.extract_text() or ""
                        text_parts.append(page_text)
                logger.info(f"Extracted {num_pages} pages via pdfplumber")
            except ImportError:
                try:
                    from PyPDF2 import PdfReader
                    reader = PdfReader(io.BytesIO(content))
                    num_pages = len(reader.pages)
                    for page in reader.pages:
                        page_text = page.extract_text() or ""
                        text_parts.append(page_text)
                    logger.info(f"Extracted {num_pages} pages via PyPDF2")
                except ImportError:
                    raise RuntimeError("No PDF library available (try: pip install PyMuPDF pdfplumber PyPDF2)")

        full_text = "\n\n".join(text_parts)
        word_count = len(full_text.split())
        return ParsedDocument(full_text.strip(), {
            "format": "pdf",
            "pages": num_pages,
            "word_count": word_count,
        })


class MarkdownParser(BaseParser):
    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            text = content.decode("utf-8", errors="replace")
            word_count = len(text.split())
            return ParsedDocument(text, {
                "format": "markdown",
                "word_count": word_count,
            })
        except Exception as e:
            logger.error(f"Markdown parsing failed: {e}")
            return ParsedDocument("", {"format": "markdown", "error": str(e)})


class TXTParser(BaseParser):
    """Plain-text parser — covers .txt (Must per spec)."""

    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            # Try utf-8 first, fall back to latin-1
            try:
                text = content.decode("utf-8")
            except UnicodeDecodeError:
                text = content.decode("latin-1", errors="replace")
            word_count = len(text.split())
            return ParsedDocument(text, {
                "format": "text",
                "word_count": word_count,
            })
        except Exception as e:
            logger.error(f"TXT parsing failed: {e}")
            return ParsedDocument("", {"format": "text", "error": str(e)})


class CSVParser(BaseParser):
    """CSV parser — header inference + row text extraction (Must spreadsheet)."""

    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            return await asyncio.get_event_loop().run_in_executor(None, self._parse_sync, content)
        except Exception as e:
            logger.error(f"CSV parsing failed: {e}")
            return ParsedDocument(f"CSV parsing error: {e}", {"format": "csv", "error": str(e)})

    def _parse_sync(self, content: bytes) -> ParsedDocument:
        try:
            text = content.decode("utf-8", errors="replace")
        except Exception:
            text = content.decode("latin-1", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        if not rows:
            return ParsedDocument("", {"format": "csv", "word_count": 0, "rows": 0})
        header = rows[0] if rows else []
        # Build searchable text: header + each row joined by " | "
        text_parts = [" | ".join(header)]
        for row in rows[1:]:
            text_parts.append(" | ".join(row))
        full_text = "\n".join(text_parts)
        word_count = len(full_text.split())
        return ParsedDocument(full_text.strip(), {
            "format": "csv",
            "word_count": word_count,
            "rows": len(rows),
            "columns": len(header),
            "header": header,
        })


class XLSXParser(BaseParser):
    """XLSX parser — via openpyxl if available, with graceful fallback."""

    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            return await asyncio.get_event_loop().run_in_executor(None, self._parse_sync, content)
        except Exception as e:
            logger.error(f"XLSX parsing failed: {e}")
            return ParsedDocument(f"XLSX parsing error: {e}", {"format": "xlsx", "error": str(e)})

    def _parse_sync(self, content: bytes) -> ParsedDocument:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True, read_only=True)
            text_parts = []
            sheet_count = len(wb.sheetnames)
            for ws in wb.worksheets:
                text_parts.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    # Filter None, join
                    cells = [str(c) if c is not None else "" for c in row]
                    # Skip entirely empty rows
                    if any(c.strip() for c in cells):
                        text_parts.append(" | ".join(cells))
            wb.close()
            full_text = "\n".join(text_parts)
            word_count = len(full_text.split())
            return ParsedDocument(full_text.strip(), {
                "format": "xlsx",
                "word_count": word_count,
                "sheets": sheet_count,
            })
        except ImportError:
            logger.warning("openpyxl not installed — falling back to CSV-like decode for XLSX")
            # Fallback: try to decode as text for basic searchability
            try:
                text = content.decode("utf-8", errors="replace")
                word_count = len(text.split())
                return ParsedDocument(text.strip(), {
                    "format": "xlsx",
                    "word_count": word_count,
                    "warning": "openpyxl not available, raw decode fallback",
                })
            except Exception as e:
                raise RuntimeError("openpyxl not available (try: pip install openpyxl)") from e


class PPTXParser(BaseParser):
    """PPTX parser — via python-pptx if available, with graceful fallback."""

    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            return await asyncio.get_event_loop().run_in_executor(None, self._parse_sync, content)
        except Exception as e:
            logger.error(f"PPTX parsing failed: {e}")
            return ParsedDocument(f"PPTX parsing error: {e}", {"format": "pptx", "error": str(e)})

    def _parse_sync(self, content: bytes) -> ParsedDocument:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            slide_count = len(prs.slides)
            for idx, slide in enumerate(prs.slides, start=1):
                text_parts.append(f"# Slide {idx}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            txt = paragraph.text.strip()
                            if txt:
                                text_parts.append(txt)
                    # Tables inside shapes
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            if any(row_text):
                                text_parts.append(" | ".join(row_text))
            full_text = "\n".join(text_parts)
            word_count = len(full_text.split())
            return ParsedDocument(full_text.strip(), {
                "format": "pptx",
                "word_count": word_count,
                "slides": slide_count,
            })
        except ImportError:
            logger.warning("python-pptx not installed — fallback raw decode for PPTX")
            try:
                text = content.decode("utf-8", errors="replace")
                word_count = len(text.split())
                return ParsedDocument(text.strip(), {
                    "format": "pptx",
                    "word_count": word_count,
                    "warning": "python-pptx not available, raw decode fallback",
                })
            except Exception as e:
                raise RuntimeError("python-pptx not available (try: pip install python-pptx)") from e


class DOCXParser(BaseParser):
    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            return await asyncio.get_event_loop().run_in_executor(None, self._parse_sync, content)
        except Exception as e:
            logger.error(f"DOCX parsing failed: {e}")
            return ParsedDocument(f"DOCX parsing error: {e}", {"format": "docx", "error": str(e)})

    def _parse_sync(self, content: bytes) -> ParsedDocument:
        try:
            import docx
            doc = docx.Document(io.BytesIO(content))
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]

            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    text_parts.append(" | ".join(row_text))

            full_text = "\n".join(text_parts)
            word_count = len(full_text.split())
            return ParsedDocument(full_text.strip(), {
                "format": "docx",
                "pages": len(doc.sections),
                "word_count": word_count,
            })
        except ImportError:
            raise RuntimeError("python-docx library is not available (try: pip install python-docx)")


class ImageParser(BaseParser):
    async def parse(self, content: bytes) -> ParsedDocument:
        try:
            return await asyncio.get_event_loop().run_in_executor(None, self._parse_sync, content)
        except Exception as e:
            logger.error(f"Image parsing failed: {e}")
            return ParsedDocument(
                content="Image OCR failed",
                metadata={"format": "image", "error": str(e), "needs_review": True},
            )

    def _parse_sync(self, content: bytes) -> ParsedDocument:
        try:
            import pytesseract
            from PIL import Image
            image = Image.open(io.BytesIO(content))
            # Try to compute confidence via image_to_data; fall back to 0.75 if unavailable
            try:
                data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
                confs = [int(c) for c in data.get("conf", []) if int(c) != -1]
                if confs:
                    confidence = sum(confs) / len(confs) / 100.0
                else:
                    confidence = 0.75
                ocr_text = pytesseract.image_to_string(image)
            except Exception:
                # Fallback path if image_to_data not available / fails
                ocr_text = pytesseract.image_to_string(image)
                confidence = 0.75
            # Clamp
            confidence = max(0.0, min(1.0, confidence))
            word_count = len(ocr_text.split())
            return ParsedDocument(
                content=ocr_text.strip(),
                metadata={
                    "format": "image",
                    "ocr_confidence": round(confidence, 3),
                    "needs_review": confidence < 0.75,
                    "word_count": word_count,
                },
            )
        except ImportError:
            logger.warning("pytesseract or Pillow not installed — returning mock OCR result")
            return ParsedDocument(
                content="Image OCR not available (pytesseract/Pillow not installed)",
                metadata={
                    "format": "image",
                    "ocr_confidence": 0.0,
                    "needs_review": True,
                    "warning": "pytesseract or Pillow not available",
                },
            )
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return ParsedDocument(
                content="Image OCR failed",
                metadata={"format": "image", "error": str(e), "needs_review": True},
            )


# Unified whitelist — single source of truth per F-40 remediation
# Covers Must spec: PDF, DOCX, PPTX, XLSX/CSV, Markdown, TXT, images
PARSERS: dict[str, type[BaseParser]] = {
    ".pdf": PDFParser,
    ".md": MarkdownParser,
    ".markdown": MarkdownParser,
    ".docx": DOCXParser,
    ".doc": DOCXParser,
    ".txt": TXTParser,
    ".csv": CSVParser,
    ".xlsx": XLSXParser,
    ".xls": XLSXParser,
    ".pptx": PPTXParser,
    ".ppt": PPTXParser,
    ".jpg": ImageParser,
    ".jpeg": ImageParser,
    ".png": ImageParser,
    ".gif": ImageParser,
    ".webp": ImageParser,
    ".svg": TXTParser,  # SVG is XML vector — treat as text, not OCR
}


class UnsupportedFormatError(Exception):
    pass


async def parse_document(filename: str, content: bytes) -> ParsedDocument:
    ext = Path(filename).suffix.lower()
    parser_cls = PARSERS.get(ext)
    if not parser_cls:
        raise UnsupportedFormatError(f"No parser for {ext}")

    parser = parser_cls(timeout=30)
    try:
        return await asyncio.wait_for(parser.parse(content), timeout=parser.timeout)
    except TimeoutError:
        logger.error(f"Parsing timed out for {filename}")
        raise
